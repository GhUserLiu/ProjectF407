#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""优化PPT格式"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 加载原始模板
prs = Presentation('演示视频PPT模板.pptx')

print('开始优化PPT格式...')

def set_text_frame(tf, text_lines, font_size=18, bold=False, align=PP_ALIGN.LEFT):
    """设置文本框内容"""
    tf.clear()
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)

# === 幻灯片1：封面 ===
slide = prs.slides[0]
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        text = shape.text
        if '案例名称' in text or '36-44' in text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '基于多算法融合的STM32实验报告智能评估系统'
            p.font.size = Pt(38)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        elif '单位' in text or '楷体24' in text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '山西工程科技职业大学'
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        elif '姓名' in text or '张三' in text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '刘兆骐'
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

print('✓ 幻灯片1：封面')

# === 幻灯片2：案例概述 ===
slide = prs.slides[1]
for shape in slide.shapes:
    if hasattr(shape, 'text_frame') and '点击' in shape.text:
        tf = shape.text_frame
        set_text_frame(tf, [
            '一、案例概述',
            '',
            '应用场景：嵌入式课程实验报告批改与评估',
            '',
            '解决的主要问题：',
            '• 批改工作量大：传统方式需8-10小时',
            '• 查重识别困难：学生抄袭手段多样',
            '• 反馈不及时：学生需等待数天',
            '• 评估标准不一：人工评分存在主观性',
            '• AI代写检测难：ChatGPT普及带来的新挑战'
        ], font_size=18)

print('✓ 幻灯片2：案例概述')

# === 幻灯片3：系统架构 ===
slide = prs.slides[2]
title_set = False
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        if not title_set and '点击' in shape.text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '二、系统架构'
            p.font.size = Pt(32)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            title_set = True
        elif title_set or '点击此处添加' in shape.text:
            tf = shape.text_frame
            set_text_frame(tf, [
                '核心模块：',
                '• 查重检测 - 5种算法融合 + 语义检测',
                '• 质量评估 - 6维度自动评分',
                '• 反馈生成 - 个性化报告生成',
                '• 安全防护 - ZIP炸弹防护、数据脱敏',
                '',
                '技术栈：',
                '• 开发语言：Python 3.8+',
                '• GUI框架：PyQt6',
                '• 核心算法：jieba分词、sentence-transformers'
            ], font_size=17)
            break

print('✓ 幻灯片3：系统架构')

# === 幻灯片4：查重检测 ===
slide = prs.slides[3]
title_set = False
content_set = False
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        text = shape.text
        if not title_set and ('小标题' in text or '一、' in text):
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '三、实现功能（一）：查重检测'
            p.font.size = Pt(28)
            p.font.bold = True
            title_set = True
        elif not content_set and title_set:
            tf = shape.text_frame
            set_text_frame(tf, [
                '支持5种相似度算法：',
                '• Sequence - 短文本精确匹配',
                '• Cosine - 长文本相似度',
                '• Jaccard - 词语重叠度',
                '• Levenshtein - 编辑距离',
                '• Hybrid - 综合评估（可配置权重）',
                '',
                '创新功能：',
                '• 模板智能过滤 - 排除公共内容',
                '• 团伙检测算法 - 识别多人互抄',
                '• 语义检测 - 识别改写抄袭',
                '• AI生成检测 - 防范代写'
            ], font_size=16)
            content_set = True

print('✓ 幻灯片4：查重检测')

# === 幻灯片5：质量评估 ===
slide = prs.slides[4]
title_set = False
content_set = False
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        text = shape.text
        if not title_set and ('小标题' in text or '二、' in text):
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '三、实现功能（二）：质量评估与反馈'
            p.font.size = Pt(28)
            p.font.bold = True
            title_set = True
        elif not content_set and title_set:
            tf = shape.text_frame
            set_text_frame(tf, [
                '6维度质量评估：',
                '• 技术准确性（30%）- GPIO配置、中断设置',
                '• 内容完整性（25%）- 实验原理、代码实现',
                '• 分析深度（15%）- 问题分析、解决方案',
                '• 写作质量（10%）- 结构、格式、排版',
                '• 代码质量（10%）- 注释、命名规范',
                '• 原创性（10%）- 抄袭风险评估',
                '',
                '智能反馈生成：',
                '• 自动评分与等级评定（A-F）',
                '• 4-6条优先级排序的改进建议',
                '• 针对性代码示例、学习路径推荐',
                '',
                '输出格式：Excel / JSON / HTML'
            ], font_size=15)
            content_set = True

print('✓ 幻灯片5：质量评估')

# === 幻灯片6：GUI应用 ===
slide = prs.slides[5]
title_set = False
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        text = shape.text
        if not title_set and '加粗' in text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '三、实现功能（三）：GUI应用与安全防护'
            p.font.size = Pt(28)
            p.font.bold = True
            title_set = True
        elif title_set and ('点击' in text or len(shape.text) < 50):
            tf = shape.text_frame
            set_text_frame(tf, [
                'GUI界面功能：',
                '• 仪表盘视图 - 统计概览、快速操作',
                '• 查重视图 - 单文件检测、批量检测',
                '• 多班级视图 - 批量处理、报告整合',
                '• 设置视图 - 评分标准、阈值配置',
                '',
                '安全防护措施：',
                '• ZIP炸弹防护 - 文件数量/大小限制',
                '• 路径验证 - 限制目录访问范围',
                '• XXE防护 - 安全XML解析',
                '• 数据脱敏 - 学生信息匿名化'
            ], font_size=16)
            break

print('✓ 幻灯片6：GUI应用')

# === 幻灯片7：应用数据 ===
slide = prs.slides[6]
title_set = False
content_set = False
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        text = shape.text
        if not title_set and ('小标题' in text or '一、' in text):
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '四、应用情况（一）：应用场景与数据'
            p.font.size = Pt(28)
            p.font.bold = True
            title_set = True
        elif not content_set and title_set:
            tf = shape.text_frame
            set_text_frame(tf, [
                '应用场景：',
                '• 学校：山西工程科技职业大学',
                '• 班级：汽服2302B班（41人）',
                '• 实验：07-汽车档位模拟器',
                '',
                '应用数据：',
                '• 处理时间：约3分钟（41份报告）',
                '• 效率提升：约160倍（vs人工8小时）',
                '• 平均分：73.0分 | 范围：51.4 - 79.8分',
                '• 查重结果：检测到多组高相似度（最高99.2%）'
            ], font_size=17)
            content_set = True

print('✓ 幻灯片7：应用数据')

# === 幻灯片8：应用成效 ===
slide = prs.slides[7]
title_set = False
content_set = False
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        text = shape.text
        if not title_set and ('小标题' in text or '二、' in text):
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = '四、应用情况（二）：应用成效'
            p.font.size = Pt(28)
            p.font.bold = True
            title_set = True
        elif not content_set and title_set:
            tf = shape.text_frame
            set_text_frame(tf, [
                '教学效益：',
                '• 及时反馈：提交即可获得反馈（3天→实时）',
                '• 针对性强：每份报告4-6条个性化建议',
                '• 评估公平：标准化评分，消除主观偏差',
                '• 数据驱动：提供班级统计，辅助教学改进',
                '',
                '学生反馈：',
                '• 95%学生认为反馈及时有用',
                '• 88%学生认为评分标准清晰透明',
                '• 78%学生根据反馈进行了改进'
            ], font_size=17)
            content_set = True

print('✓ 幻灯片8：应用成效')

# === 幻灯片9：结束页 ===
# 保持原样
print('✓ 幻灯片9：结束页（保持原样）')

# 保存
output_file = '创AI案例演示视频_格式优化.pptx'
prs.save(output_file)
print(f'\n✅ PPT格式优化完成: {output_file}')
