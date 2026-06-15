#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""填充PPT内容"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

# 加载模板
prs = Presentation('演示视频PPT模板.pptx')

print('开始填充PPT内容...')

# 幻灯片1：封面
slide = prs.slides[0]
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        if '案例名称' in shape.text or '楷体36' in shape.text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = '基于多算法融合的STM32实验报告智能评估系统'
            p.font.size = Pt(36)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        elif '单位' in shape.text or '楷体24' in shape.text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = '山西工程科技职业大学'
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        elif '姓名' in shape.text or '张三' in shape.text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = '刘兆骐'
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

print('幻灯片1：封面 - 完成')

# 幻灯片2：案例概述
slide = prs.slides[1]
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        if '点击此处添加' in shape.text:
            shape.text_frame.clear()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = '一、案例概述'
            p.font.size = Pt(28)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # 添加内容
            text_lines = [
                '',
                '应用场景：嵌入式课程实验报告批改与评估',
                '',
                '解决的主要问题：',
                '• 批改工作量大：传统方式需8-10小时',
                '• 查重识别困难：学生抄袭手段多样',
                '• 反馈不及时：学生需等待数天',
                '• 评估标准不一：人工评分存在主观性',
                '• AI代写检测难：ChatGPT普及带来的新挑战'
            ]

            for line in text_lines:
                if line:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(18)
                    p.alignment = PP_ALIGN.LEFT

print('幻灯片2：案例概述 - 完成')

# 幻灯片3：系统架构
slide = prs.slides[2]
content_added = False
for shape in slide.shapes:
    if hasattr(shape, 'text_frame') and not content_added:
        if '点击此处添加' in shape.text:
            # 找到主标题位置
            shape.text_frame.clear()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = '二、系统架构'
            p.font.size = Pt(28)
            p.font.bold = True

            content_added = True
    elif hasattr(shape, 'text_frame') and content_added:
        # 内容区域
        shape.text_frame.clear()
        tf = shape.text_frame
        tf.word_wrap = True

        text_lines = [
            '核心模块：',
            '• 查重检测 - 5种算法融合 + 语义检测',
            '• 质量评估 - 6维度自动评分',
            '• 反馈生成 - 个性化报告生成',
            '• 安全防护 - ZIP炸弹防护、数据脱敏',
            '',
            '技术栈：',
            '• 开发语言：Python 3.8+',
            '• GUI框架：PyQt6',
            '• 核心算法：jieba分词、sentence-transformers'
        ]

        for i, line in enumerate(text_lines):
            p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.space_before = Pt(4)

print('幻灯片3：系统架构 - 完成')

# 幻灯片4：查重检测
slide = prs.slides[3]
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        if '一、点击此处添加小标题' in shape.text or '一、' in shape.text:
            shape.text_frame.clear()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = '三、实现功能（一）：查重检测'
            p.font.size = Pt(26)
            p.font.bold = True
        elif '实现功能' in shape.text:
            shape.text_frame.clear()
            tf = shape.text_frame

            text_lines = [
                '支持5种相似度算法：',
                '• Sequence - 短文本精确匹配',
                '• Cosine - 长文本相似度',
                '• Jaccard - 词语重叠度',
                '• Levenshtein - 编辑距离',
                '• Hybrid - 综合评估（可配置权重）',
                '',
                '创新功能：',
                '• 模板智能过滤 - 排除公共内容',
                '• 团伙检测算法 - 识别多人互抄',
                '• 语义检测 - 识别改写抄袭',
                '• AI生成检测 - 防范代写'
            ]

            for i, line in enumerate(text_lines):
                if i < len(tf.paragraphs):
                    p = tf.paragraphs[i]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(15)
                p.space_before = Pt(3)

print('幻灯片4：查重检测 - 完成')

# 幻灯片5：质量评估
slide = prs.slides[4]
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        if '二、点击此处添加小标题' in shape.text or '二、' in shape.text:
            shape.text_frame.clear()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = '三、实现功能（二）：质量评估与反馈'
            p.font.size = Pt(26)
            p.font.bold = True
        elif '点击此处添加' in shape.text and '实现功能' not in shape.text:
            shape.text_frame.clear()
            tf = shape.text_frame

            text_lines = [
                '6维度质量评估：',
                '• 技术准确性（30%） - GPIO配置、中断设置',
                '• 内容完整性（25%） - 实验原理、代码实现',
                '• 分析深度（15%） - 问题分析、解决方案',
                '• 写作质量（10%） - 结构、格式、排版',
                '• 代码质量（10%） - 注释、命名规范',
                '• 原创性（10%） - 抄袭风险评估',
                '',
                '智能反馈生成：',
                '• 自动评分与等级评定（A-F）',
                '• 4-6条优先级排序的改进建议',
                '• 针对性代码示例、学习路径推荐',
                '',
                '输出格式：Excel / JSON / HTML'
            ]

            for i, line in enumerate(text_lines):
                if i < len(tf.paragraphs):
                    p = tf.paragraphs[i]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.space_before = Pt(2)

print('幻灯片5：质量评估 - 完成')

# 幻灯片6：GUI应用
slide = prs.slides[5]
content_added = False
for shape in slide.shapes:
    if hasattr(shape, 'text_frame') and not content_added:
        if '点击此处添加' in shape.text and '加粗' in shape.text:
            shape.text_frame.clear()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = '三、实现功能（三）：GUI应用与安全防护'
            p.font.size = Pt(26)
            p.font.bold = True
            content_added = True
    elif hasattr(shape, 'text_frame') and content_added:
        shape.text_frame.clear()
        tf = shape.text_frame

        text_lines = [
            'GUI界面功能：',
            '• 仪表盘视图 - 统计概览、快速操作',
            '• 查重视图 - 单文件检测、批量检测',
            '• 多班级视图 - 批量处理、报告整合',
            '• 设置视图 - 评分标准、阈值配置',
            '',
            '安全防护措施：',
            '• ZIP炸弹防护 - 文件数量/大小限制',
            '• 路径验证 - 限制目录访问范围',
            '• XXE防护 - 安全XML解析',
            '• 数据脱敏 - 学生信息匿名化'
        ]

        for i, line in enumerate(text_lines):
            if i < len(tf.paragraphs):
                p = tf.paragraphs[i]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(15)
            p.space_before = Pt(3)

print('幻灯片6：GUI应用 - 完成')

# 幻灯片7：应用场景与数据
slide = prs.slides[6]
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        if '一、点击此处添加小标题' in shape.text or '一、' in shape.text:
            shape.text_frame.clear()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = '四、应用情况（一）：应用场景与数据'
            p.font.size = Pt(26)
            p.font.bold = True
        elif '应用情况' in shape.text:
            shape.text_frame.clear()
            tf = shape.text_frame

            text_lines = [
                '应用场景：',
                '• 学校：山西工程科技职业大学',
                '• 班级：汽服2302B班（41人）',
                '• 实验：07-汽车档位模拟器',
                '',
                '应用数据：',
                '• 处理时间：约3分钟（41份报告）',
                '• 效率提升：约160倍（vs人工8小时）',
                '• 平均分：73.0分 | 范围：51.4 - 79.8分',
                '• 查重结果：检测到多组高相似度报告（最高99.2%）'
            ]

            for i, line in enumerate(text_lines):
                if i < len(tf.paragraphs):
                    p = tf.paragraphs[i]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(15)
                p.space_before = Pt(3)

print('幻灯片7：应用数据 - 完成')

# 幻灯片8：应用成效
slide = prs.slides[7]
for shape in slide.shapes:
    if hasattr(shape, 'text_frame'):
        if '二、点击此处添加小标题' in shape.text or '二、' in shape.text:
            shape.text_frame.clear()
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = '四、应用情况（二）：应用成效'
            p.font.size = Pt(26)
            p.font.bold = True
        elif '点击此处添加' in shape.text and '应用情况' not in shape.text:
            shape.text_frame.clear()
            tf = shape.text_frame

            text_lines = [
                '教学效益：',
                '• 及时反馈：提交即可获得反馈（3天→实时）',
                '• 针对性强：每份报告4-6条个性化建议',
                '• 评估公平：标准化评分，消除主观偏差',
                '• 数据驱动：提供班级统计，辅助教学改进',
                '',
                '学生反馈：',
                '• 95%学生认为反馈及时有用',
                '• 88%学生认为评分标准清晰透明',
                '• 78%学生根据反馈进行了改进'
            ]

            for i, line in enumerate(text_lines):
                if i < len(tf.paragraphs):
                    p = tf.paragraphs[i]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(15)
                p.space_before = Pt(3)

print('幻灯片8：应用成效 - 完成')

# 幻灯片9：结束页保持不变
print('幻灯片9：结束页 - 保持原样')

# 保存
output_file = '创AI案例演示视频.pptx'
prs.save(output_file)
print(f'\n✓ PPT生成成功: {output_file}')
