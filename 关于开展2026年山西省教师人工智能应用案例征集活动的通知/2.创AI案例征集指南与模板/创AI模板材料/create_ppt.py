#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""创建全新的PPT演示文稿"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建新的演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

def add_title_slide(prs, title, subtitle):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
    slide = prs.slides.add_slide(slide_layout)

    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title

    # 设置副标题
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle

    return slide

def add_content_slide(prs, title, content_lines, font_size=18):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)

    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title

    # 设置内容
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = line
        p.font.size = Pt(font_size)
        p.space_after = Pt(10)

        # 如果是列表项（以•开头）
        if line.strip().startswith('•'):
            p.level = 0
        # 如果是二级列表
        elif line.strip().startswith('  '):
            p.level = 1

    return slide

def add_two_content_slide(prs, title, left_content, right_content, font_size=16):
    """添加双栏内容幻灯片"""
    # 使用空白布局，手动创建文本框
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 添加标题
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 左栏内容
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(4.2)
    height = Inches(4)
    left_box = slide.shapes.add_textbox(left, top, width, height)
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, line in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.space_after = Pt(6)

    # 右栏内容
    left = Inches(5.3)
    right_box = slide.shapes.add_textbox(left, top, width, height)
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, line in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.space_after = Pt(6)

    return slide

print('开始创建PPT...')

# === 幻灯片1：封面 ===
add_title_slide(prs,
    '基于多算法融合的STM32\n实验报告智能评估系统',
    '山西工程科技职业大学  刘兆骐')
print('✓ 幻灯片1：封面')

# === 幻灯片2：案例概述 ===
add_content_slide(prs, '一、案例概述', [
    '应用场景：嵌入式课程实验报告批改与评估',
    '',
    '解决的主要问题：',
    '• 批改工作量大：传统方式需8-10小时',
    '• 查重识别困难：学生抄袭手段多样',
    '• 反馈不及时：学生需等待数天',
    '• 评估标准不一：人工评分存在主观性',
    '• AI代写检测难：ChatGPT普及带来的新挑战'
], font_size=18)
print('✓ 幻灯片2：案例概述')

# === 幻灯片3：系统架构 ===
add_content_slide(prs, '二、系统架构', [
    '核心模块：',
    '• 查重检测 - 5种算法融合 + 语义检测',
    '• 质量评估 - 6维度自动评分',
    '• 反馈生成 - 个性化报告生成',
    '• 安全防护 - ZIP炸弹防护、数据脱敏',
    '',
    '技术栈：',
    '• 开发语言：Python 3.8+',
    '• GUI框架：PyQt6',
    '• 核心算法：jieba分词、sentence-transformers'
], font_size=17)
print('✓ 幻灯片3：系统架构')

# === 幻灯片4：查重检测 ===
add_two_content_slide(prs, '三、实现功能（一）：查重检测',
    ['支持5种相似度算法：',
     '• Sequence - 短文本精确匹配',
     '• Cosine - 长文本相似度',
     '• Jaccard - 词语重叠度',
     '• Levenshtein - 编辑距离',
     '• Hybrid - 综合评估'],
    ['创新功能：',
     '• 模板智能过滤',
     '  排除公共内容',
     '• 团伙检测算法',
     '  识别多人互抄',
     '• 语义检测',
     '  识别改写抄袭',
     '• AI生成检测',
     '  防范代写'],
    font_size=15)
print('✓ 幻灯片4：查重检测')

# === 幻灯片5：质量评估 ===
add_two_content_slide(prs, '三、实现功能（二）：质量评估与反馈',
    ['6维度质量评估：',
     '• 技术准确性（30%）',
     '• 内容完整性（25%）',
     '• 分析深度（15%）',
     '• 写作质量（10%）',
     '• 代码质量（10%）',
     '• 原创性（10%）'],
    ['智能反馈生成：',
     '• 自动评分与等级评定（A-F）',
     '• 4-6条优先级排序建议',
     '• 针对性代码示例',
     '• 学习路径推荐',
     '',
     '输出格式：',
     'Excel / JSON / HTML'],
    font_size=14)
print('✓ 幻灯片5：质量评估')

# === 幻灯片6：GUI应用 ===
add_two_content_slide(prs, '三、实现功能（三）：GUI应用与安全防护',
    ['GUI界面功能：',
     '• 仪表盘视图',
     '  统计概览、快速操作',
     '• 查重视图',
     '  单文件检测、批量检测',
     '• 多班级视图',
     '  批量处理、报告整合',
     '• 设置视图',
     '  评分标准、阈值配置'],
    ['安全防护措施：',
     '• ZIP炸弹防护',
     '  文件数量/大小限制',
     '• 路径验证',
     '  限制目录访问范围',
     '• XXE防护',
     '  安全XML解析',
     '• 数据脱敏',
     '  学生信息匿名化'],
    font_size=14)
print('✓ 幻灯片6：GUI应用')

# === 幻灯片7：应用数据 ===
add_content_slide(prs, '四、应用情况（一）：应用场景与数据', [
    '应用场景：',
    '• 学校：山西工程科技职业大学',
    '• 班级：汽服2302B班（41人）',
    '• 实验：07-汽车档位模拟器',
    '',
    '应用数据：',
    '• 处理时间：约3分钟（41份报告）',
    '• 效率提升：约160倍（vs人工8小时）',
    '• 平均分：73.0分 | 范围：51.4 - 79.8分',
    '• 查重结果：检测到多组高相似度（最高99.2%）'
], font_size=18)
print('✓ 幻灯片7：应用数据')

# === 幻灯片8：应用成效 ===
add_two_content_slide(prs, '四、应用情况（二）：应用成效',
    ['教学效益：',
     '• 及时反馈',
     '  3天→实时',
     '• 针对性强',
     '  4-6条个性化建议',
     '• 评估公平',
     '  标准化评分',
     '• 数据驱动',
     '  班级统计分析'],
    ['学生反馈：',
     '• 95%认为反馈及时有用',
     '• 88%认为标准清晰透明',
     '• 78%根据反馈改进',
     '',
     '',
     '',
     '',
     ''],
    font_size=16)
print('✓ 幻灯片8：应用成效')

# === 幻灯片9：结束页 ===
add_title_slide(prs, '谢谢！', '')
print('✓ 幻灯片9：结束页')

# 保存
output_file = '创AI案例演示视频_新版.pptx'
prs.save(output_file)

print(f'\n✅ PPT创建完成: {output_file}')
print(f'总页数: {len(prs.slides)}')
